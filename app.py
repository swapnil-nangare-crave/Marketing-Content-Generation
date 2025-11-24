
# Unified AI Content Generator (Interactive Streamlit Layout)

import os
import re
import html
import warnings
import requests #type: ignore
import streamlit as st  # type: ignore

# Document reading
from PyPDF2 import PdfReader  # type: ignore
from docx import Document as DocxDocument  # type: ignore
from pptx import Presentation  # type: ignore

# HANA + LangChain + Azure clients
from hdbcli import dbapi  # type: ignore
from langchain_community.vectorstores.hanavector import HanaDB  # type: ignore
from langchain_openai import AzureOpenAIEmbeddings  # type: ignore
from openai import AzureOpenAI  # type: ignore

warnings.filterwarnings("ignore", category=DeprecationWarning)

# Perplexity config (if used)

PERPLEXITY_API_KEY = st.secrets.get("perplexity", {}).get("api_key", "")
PERPLEXITY_API_URL = st.secrets.get("perplexity", {}).get("api_url", "https://api.perplexity.ai/search")


def perplexity_search(query, max_results=5):
    """Fetch results from Perplexity.ai"""
    if not PERPLEXITY_API_KEY:
        return ""
    payload = {"query": query}
    headers = {
        "Authorization": f"Bearer {PERPLEXITY_API_KEY}",
        "Content-Type": "application/json",
    }
    try:
        response = requests.post(PERPLEXITY_API_URL, json=payload, headers=headers, timeout=15)
        response.raise_for_status()
        data = response.json()
        results = []
        if "answer" in data:
            results.append(data["answer"])
        elif "data" in data and isinstance(data["data"], list):
            for item in data["data"][:max_results]:
                if "text" in item:
                    results.append(item["text"])
        return "\n".join(results)
    except Exception as e:
        return f"Perplexity API Error: {e}"


# File / URL text extraction

def extract_text_from_file(file):
    """Extract text from TXT / PDF / DOCX / PPTX uploads"""
    text = ""
    name = file.name.lower()

    try:
        if name.endswith(".txt"):
            text = file.read().decode("utf-8", errors="ignore")

        elif name.endswith(".pdf"):
            pdf = PdfReader(file)
            for page in pdf.pages:
                text += page.extract_text() or ""

        elif name.endswith(".docx"):
            doc = DocxDocument(file)
            text = "\n".join([p.text for p in doc.paragraphs])

        elif name.endswith(".pptx"):
            ppt = Presentation(file)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
    except Exception:
        # best-effort extraction; avoid breaking generation flow
        pass

    return text.strip()


def extract_text_from_url(url):
    """Fetch and extract readable text from a given URL using Perplexity (fallback)"""
    if not url or not url.strip():
        return ""
    try:
        if PERPLEXITY_API_KEY:
            payload = {"query": f"Extract main article content from: {url}"}
            headers = {
                "Authorization": f"Bearer {PERPLEXITY_API_KEY}",
                "Content-Type": "application/json",
            }
            response = requests.post(PERPLEXITY_API_URL, json=payload, headers=headers, timeout=20)
            response.raise_for_status()
            data = response.json()
            if "answer" in data:
                return data["answer"]
            elif "data" in data and isinstance(data["data"], list):
                return "\n".join([d.get("text", "") for d in data["data"]])
        # fallback: simple HTTP GET (less reliable for clean article text)
        r = requests.get(url, timeout=10)
        if r.ok:
            return r.text
    except Exception as e:
        return f"Error extracting URL content: {e}"
    return ""


import streamlit as st
import markdown

# ======================================================
# Streamlit UI Setup
# ======================================================

st.set_page_config(page_title="AI Content Hub", layout="wide")

# ======================================================
# Sidebar - Global Settings
# ======================================================

st.sidebar.markdown("## ⚙️ Content Configuration")

# Content Type Selection
content_type = st.sidebar.selectbox(
    "📄 Select Content Type",
    ["Blog", "Video Script"],
)

# Dynamic Title Logic
if content_type == "Blog":
    dynamic_title = "AI Blog Generator"
else:
    dynamic_title = "AI Video Script Generator"

tone = st.sidebar.selectbox(
    "🎨 Tone",
    [
        "Professional", "Friendly", "Authoritative", "Playful", "Inspirational",
        "Conversational", "Casual", "Semi-casual", "Business professional",
        "Approachable", "Informative", "Assertive", "Engaging",
        "Visionary (for Thought Leadership)", "Confident", "Data-driven",
        "Plainspoken / Direct", "Witty", "Storytelling"
    ],
)

target_audience = st.sidebar.selectbox(
    "🎯 Target Audience",
    ["Senior Management", "Middle Management", "Junior / Entry Level Staff"],
)

industry = st.sidebar.text_input(
    "🏢 Industry (optional)",
    placeholder="e.g., Manufacturing, Retail, Technology",
)

# Blog Word Limit OR Script Timing
if content_type == "Blog":
    word_limit = st.sidebar.slider("📝 Word Limit", 300, 2000, 1000, step=100)
    time_limit = None
else:
    time_limit = st.sidebar.slider("⏱️ Video Duration (minutes)", 0.5, 10.0, 1.5, step=0.5)
    word_limit = None

cta_options = [
    "Talk to our experts", "Learn more about our solutions", "Book a free consultation",
    "Book Assessment", "Contact us today", "Download the full guide", "Request a demo",
]
cta_choice = st.sidebar.selectbox("📢 Call-to-Action (CTA)", cta_options)

# ======================================================
# Custom Dynamic HTML Header
# ======================================================

st.markdown(f"""
<style>
.block-container {{
    padding-top: 0rem !important;
}}

.custom-title-container {{
    display: flex;
    align-items: center;
    gap: 12px;
    margin-top: -40px;
    margin-bottom: -20px;
}}

.custom-title-icon {{
    width: 48px;
    height: 48px;
}}

.custom-title-text {{
    font-size: 42px;
    font-weight: 700;
}}

.custom-subtitle {{
    margin-top: -10px;
    font-size: 16px;
    color: #555;
}}
</style>

<div class="custom-title-container">
    <img src="https://em-content.zobj.net/source/microsoft-teams/363/brain_1f9e0.png" class="custom-title-icon">
    <div class="custom-title-text">{dynamic_title}</div>
</div>

<div class="custom-subtitle">
    AI-powered content creation for all your marketing needs
</div>
""", unsafe_allow_html=True)

# Streamlit Title (Dynamic Too)
st.title(f"🧠 {dynamic_title}")
st.markdown("AI-powered content creation for all your marketing needs")
st.markdown("<br><br>", unsafe_allow_html=True)



# ======================================================
# Session State
# ======================================================

if "output" not in st.session_state:
    st.session_state.output = ""

if "seo_results" not in st.session_state:
    st.session_state.seo_results = {}

# ======================================================
# Upload Section + Reference URLs
# ======================================================

col1, col2 = st.columns([1.2, 1.8])

with col1:
    st.markdown("#### 📎 Upload Documents")
    uploaded_files = st.file_uploader(
        "Upload Reference Files (TXT, PDF, DOCX, PPTX)",
        type=["txt", "pdf", "docx", "pptx"],
        accept_multiple_files=True,
        label_visibility="collapsed",
    )
    if uploaded_files:
        for f in uploaded_files:
            st.markdown(f"- 📄 *{f.name}*")

with col2:
    st.markdown("#### 🔗 Reference URLs")
    reference_urls = st.text_area(
        "Add Reference URLs (comma-separated)",
        placeholder="https://example.com/page1, https://example.com/page2",
        height=70,
        label_visibility="collapsed"
    )
    url_list = [url.strip() for url in reference_urls.split(",") if url.strip()]

#st.subheader("💬 Prompt / Inputs")
query = st.text_input("**Enter your topic:**")

#st.markdown("<hr>", unsafe_allow_html=True)
# ======================================================
# Main Layout — Left Inputs / Right Output
# ======================================================

left, right = st.columns([1, 2])

with left:
    

    # SEO settings for blog
    if content_type == "Blog":
        st.subheader("🔎 SEO Settings")
        primary_keyword = st.text_input("Primary Keyword")
        lsi_keywords_input = st.text_input("LSI / Variations (comma-separated)")
        lsi_keywords = [k.strip() for k in lsi_keywords_input.split(",") if k.strip()]
    else:
        primary_keyword = ""
        lsi_keywords = []

    #st.markdown("<hr style='margin:0; border:0.5px solid #e0e0e0;'>", unsafe_allow_html=True)

    # Generate Button
    generate_button = st.button(f"Generate {content_type}")

    # Refinement Section
    st.markdown("**Refine / Edit generated output**")
    refine_instruction = st.text_area(
        "Enter refinement instruction (e.g., make tone more formal, shorten intro):",
        height=180
    )
    apply_refine = st.button("Apply Changes")

    st.markdown("<hr style='margin:0; border:0.5px solid #e0e0e0;'>", unsafe_allow_html=True)

    if st.button("Clear Output"):
        st.session_state.output = ""
        st.session_state.seo_results = {}
        st.success("Output cleared.")

# ======================================================
# Output Section
# ======================================================

with right:
    st.markdown("### 📝 Output")

    st.markdown("""
        <style>
        .output-box {
            border: 2px solid #E0E0E0;
            border-radius: 12px;
            padding: 20px;
            background-color: #F9FAFB;
            height: 100vh;
            box-shadow: 0 1px 3px rgba(0,0,0,0.05);
            overflow-y: auto;
        }
        </style>
    """, unsafe_allow_html=True)

    if st.session_state.output:
        clean_output = (
            st.session_state.output
            .replace("```markdown", "")
            .replace("```", "")
            .lstrip()
        )

        html_content = markdown.markdown(clean_output)

        st.markdown(f"""
        <div class='output-box'>
            {html_content}
        </div>
        """, unsafe_allow_html=True)

    else:
        st.markdown("<div class='output-box'><em>Generated output will appear here.</em></div>", unsafe_allow_html=True)


# Services initialization (HANA + Azure)

def init_services():
    # Check for secrets
    if "database" in st.secrets and "address" in st.secrets["database"]:
        st.write("Database secrets found.")
    else:
        st.error("Database secrets not found.")
        raise Exception("Database secrets not found.")

    if "azure" in st.secrets and "openai_endpoint" in st.secrets["azure"]:
        st.write("Azure secrets found.")
    else:
        st.error("Azure secrets not found.")
        raise Exception("Azure secrets not found.")

    try:
        connection = dbapi.connect(
            address=st.secrets["database"]["address"],
            port=st.secrets["database"]["port"],
            user=st.secrets["database"]["user"],
            password=st.secrets["database"]["password"],
            encrypt=True,
            autocommit=True,
            sslValidateCertificate=False,
        )
    except Exception as e:
        st.error(f"HANA connection error: {e}")
        raise

    try:
        client = AzureOpenAI(
            azure_endpoint=st.secrets["azure"]["openai_endpoint"],
            api_key=st.secrets["azure"]["api_key"],
            api_version=st.secrets["azure"]["api_version"],
        )
    except Exception as e:
        st.error(f"Azure OpenAI client error: {e}")
        raise

    try:
        embeddings = AzureOpenAIEmbeddings(
            azure_deployment=st.secrets["azure"]["embeddings_deployment"],
            openai_api_version=st.secrets["azure"]["embeddings_api_version"],
            api_key=st.secrets["azure"]["api_key"],
            azure_endpoint=st.secrets["azure"]["openai_endpoint"],
        )
    except Exception as e:
        st.error(f"Azure Embeddings client error: {e}")
        raise

    try:
        db = HanaDB(
            embedding=embeddings,
            connection=connection,
            table_name="MARKETING_APP_CONTENT_GENERATION",
        )
    except Exception as e:
        st.error(f"HanaDB init error: {e}")
        raise

    return db, client

# Retrieval / RAG logic (fixed hana_text initialization)

def retrieve_content(query, uploaded_files, url_list, db):
    """Retrieve content from uploaded files, URLs, HANA, or Perplexity fallback"""
    # 1) Uploaded files
    uploaded_text = ""
    if uploaded_files:
        for f in uploaded_files:
            uploaded_text += extract_text_from_file(f) + "\n"
    if uploaded_text.strip():
        return uploaded_text.strip()

    # 2) URLs
    url_text = ""
    if url_list:
        for url in url_list:
            url_text += extract_text_from_url(url) + "\n"
    if url_text.strip():
        return url_text.strip()

    # 3) HANA similarity search
    hana_text = ""
    try:
        docs = db.similarity_search(query, k=20) if db is not None else []
        # docs may be list of Document-like objects
        if docs:
            hana_text = "\n".join([getattr(doc, "page_content", str(doc)) for doc in docs])
    except Exception:
        # ensure hana_text remains defined (avoid UnboundLocalError)
        hana_text = hana_text or ""
    if hana_text.strip():
        return hana_text.strip()

    # 4) Fallback: Perplexity search of the query
    try:
        p_text = perplexity_search(query)
        if p_text and not p_text.lower().startswith("perplexity api error"):
            return p_text
    except Exception:
        pass

    return ""

# Prompt generators (your original, slightly condensed)

def generate_prompt_guidelines(tone, target_audience):
    tone_guidelines = {
        "Professional": "Use clear, concise, and confident language. Focus on credibility, precision, and business relevance.",
        "Friendly": "Use warm, conversational, and easy-to-understand language. Maintain professionalism but sound approachable.",
        "Authoritative": "Use confident, expert-driven language. Provide strong arguments and data-backed insights.",
        "Playful": "Use witty, light-hearted, and creative phrasing. Keep the tone fun yet informative, with clever transitions.",
        "Inspirational": "Use motivational and uplifting language. Focus on positive change, growth, and vision-driven storytelling.",
        "Conversational": "Write as if speaking naturally to the reader. Use a relaxed and engaging tone with simple, flowing sentences.",
        "Casual": "Keep the language light, informal, and easy to follow. Avoid jargon; use contractions and relatable examples.",
        "Semi-casual": "Balance professionalism with friendliness. Use polite, natural phrasing that feels human but credible.",
        "Business professional": "Maintain a formal, respectful tone suitable for executive audiences. Emphasize clarity, accuracy, and authority.",
        "Approachable": "Use inclusive and welcoming language. Avoid overly technical terms; sound supportive and open to dialogue.",
        "Informative": "Focus on clarity and explanation. Use structured, factual sentences that educate the reader efficiently.",
        "Assertive": "Use confident, decisive language. Clearly express opinions or recommendations without sounding aggressive.",
        "Authoritative": "Adopt an expert voice with data-backed reasoning. Establish trust through precision and confidence.",
        "Engaging": "Use dynamic, audience-focused language. Vary rhythm, include questions, and keep readers emotionally connected.",
        "First person usage + Visionary (for Thought Leadership Articles)": "Use 'I' or 'we' statements to express personal experience and vision. Inspire forward-thinking perspectives and leadership insights.",
        "Confident": "Use strong, assured language with active voice. Present ideas as well-founded and impactful.",
        "Data-driven": "Use factual, analytical language supported by evidence and statistics. Focus on insights, accuracy, and quantifiable outcomes.",
        "Plainspoken or direct": "Be concise, honest, and transparent. Avoid buzzwords and fluff; prioritize clarity and directness.",
        "Witty (a bit of humour for special cases)": "Add subtle humor or clever turns of phrase. Keep it tasteful, intelligent, and contextually relevant.",
        "Storytelling": "Use narrative-driven language with emotional pull. Build flow with characters, challenges, and resolutions to keep readers immersed."
    }
    audience_guidelines = {
        "Senior Management": (
            "Focus on strategic insights, ROI, and business impact. Use concise, high-level language. Avoid unnecessary technical details."
        ),
        "Middle Management": (
            "Provide actionable guidance, practical steps, and process-oriented insights. Balance strategic context with implementation advice."
        ),
        "Junior/Entry Level Staff": (
            "Explain clearly, use simple examples, and avoid jargon. Focus on learning, awareness, and foundational concepts."
        ),
    }
    tone_instruction = tone_guidelines.get(tone, "")
    audience_instruction = audience_guidelines.get(target_audience, "")
    return tone_instruction, audience_instruction

def enforce_word_limit(text, limit):
    """Trim text to exact word limit. Keeps first `limit` words."""
    if not limit or limit <= 0:
     return text
    words = text.split()
    if len(words) <= limit:
     return text
    trimmed = " ".join(words[:limit])
    # ensure it ends gracefully
    if not trimmed.endswith((".", "!", "?")):
     trimmed = trimmed.rstrip(',;:') + '.'
    return trimmed
       

def generate_blog_prompt(tone, target_audience, industry, query, word_limit, final_content,
                         primary_keyword, lsi_keywords, cta_text):
    tone_instruction, audience_instruction = generate_prompt_guidelines(tone, target_audience)
    
    # Strict word limits instruction
    if word_limit:
        lower = max( max(1, word_limit - 20), 1 )
        upper = word_limit + 20
        word_limit_instruction = (
        f"The final blog MUST be between {lower} and {upper} words. "
        f"Do NOT exceed this range. Stop immediately once you reach the word limit."
        )
    else:
        word_limit_instruction = ""

    return f"""
You are an experienced B2B blog writer specializing in SAP, AI, and enterprise technology domains.
Your goal is to create a marketing-grade, SEO-optimized, structured, and natural blog aligned with
enterprise communication standards. Follow these exact rules:

=====================================================
🎯 TONE & STYLE
=====================================================
- The language must always sound **natural, human, and conversational** — not robotic or AI-generated.
- Be authoritative, clear, and practical — like McKinsey insights simplified by 20%.
- Every line must add business value and flow naturally.
- Avoid over-formal phrasing, buzzwords, or filler lines (“In today’s world”, “cutting-edge”, etc.).
- Every line must add business value while flowing smoothly.
- Write like a smart consultant guiding a professional audience.
- Tone: {tone_instruction}
- Audience: {audience_instruction}

**Strict Blog Structure and Formatting Guidelines:**

1️ **Title**
   - Must be the first line of the blog.
   - The title must appear in h1 style to it of markdown.
   - Keep it short, focused, and benefit-driven.
   - Keep it under 12 words and include the primary keyword naturally.
   - Avoid clickbait or overpromises.
   - Recommended formats:
       • How [X] Helps [Y] Achieve [Z]
       • [Number] Ways to [Achieve Outcome]
       • Why [X] Isn’t Working—And How to Fix It
       • [Phrase] in the Age of [Trend]

2️ **Introduction**
   - Begin with a real scenario or insight (max 4 lines).
   - Build context in 1–2 short paragraphs.
   - Smoothly bridge to the topic—no “In this blog we’ll discuss”.

3️⃣ **Body Sections (3–5 H2s)**
   - Use descriptive or question-style subheadings (##).
    Each section must include **at least 3 rich paragraphs** that:
       • Explain the business context or challenge clearly.
       • Add examples, relatable scenarios, or short client-style stories.
       • Include supporting insights, stats, or outcomes.
       • End with a meaningful business takeaway.
   - Use connecting phrases (“This means…”, “For example…”, “In practice…”) to sound more natural.
   - Maintain a balance between data and narrative — avoid sounding like a report.
   - For shorter word limits (under 800 words), focus on **depth over quantity**: fewer sections, more substance per section.
   - Include supporting data if relevant (“According to a 2025 SAPinsider study…”).
   - Use lists only when they enhance clarity.

4️⃣ **Conclusion (Action-Oriented)**
   - Title example: “Accelerate Your Journey with [Solution]” or “Unlock the Future of [Topic]”.
   - Summarize key insights in 3–4 lines.
   - End with a strong CTA: “{cta_text}”

=====================================================
📊 SEO REQUIREMENTS
=====================================================
- Primary Keyword: "{primary_keyword}"
  • Use in Title, Intro (first 100 words), at least one H2, and 2–3 times in the body.
- LSI Keywords: {', '.join(lsi_keywords) or 'none'}
  • Include naturally where relevant, never stuffed.
- Optimize for readability and human tone — not keyword density.
- Use Markdown headings (##, ###) for structure.

=====================================================
🧱 CONTENT CONTEXT
=====================================================
Industry: {industry or "Enterprise / B2B"}
Word Limit: ~{word_limit} words
Topic: "{query}"

REFERENCE CONTENT:
{final_content or '[No reference content provided]'}

=====================================================
💡 FINAL INSTRUCTION
=====================================================
Write the blog in one coherent piece — no step-by-step notes, no bullet outlines, no commentary.
Output only the **final polished blog**, formatted for publishing.
Ensure:
- Do not add ```markdown at start and ``` at the end of response.
- Leave one blank line after the title before the introduction begins.
- The rest of the structure strictly follows the format above.
- Ensure the blog reads naturally and conversationally — it should sound like expert storytelling, not a technical report.
- When the total word limit is under 800, prioritize deeper insights per section instead of squeezing in more headings.
- Headings should be in bold 

{word_limit_instruction}

"""

def generate_video_prompt(tone, target_audience, industry, final_content,cta_text,query, time_limit):
    tone_instruction, audience_instruction = generate_prompt_guidelines(tone, target_audience)
    # Approx. 15 seconds per scene → 4 scenes per minute
    scenes = max(4, int(time_limit * 4))
    scene_duration = int((time_limit * 60) / scenes)
    return f"""
    
You are a professional **B2B video scriptwriter** who creates powerful marketing narratives.
Write a **timestamp-based video script** for the topic: "{query}" in the {industry or "enterprise"} industry.

=====================================================
🎬 STRUCTURE
=====================================================
Each scene must include:
- **Timestamp** (e.g. 0:00–0:{scene_duration:02d})
- **scene name**
- **Visuals:** On-screen visuals or camera direction
- **Narration:** Voiceover content

start each section in new line

=====================================================
🕒 TOTAL DURATION & SCENES
=====================================================
- Total video duration: ~{time_limit} minute(s)
- Divide into ~{scenes} scenes (~{scene_duration} seconds each)
- End with a personalized Call-to-Action

=====================================================
📖 STORYLINE FLOW
=====================================================
1️⃣ **Problem Introduction** – hook the viewer immediately (Scene 1)
2️⃣ **Product or Brand Introduction** – what you offer and why it matters
3️⃣ **Key Feature Highlights** – focus on impact, not just specs
4️⃣ **Benefits** – how it solves real challenges
5️⃣ **Real-Life Example or Case Study** – add credibility
6️⃣ **Call-to-Action & Closing** – must sound human, confident, and aligned with "{cta_text}"

=====================================================
🗣️ LANGUAGE & STYLE
=====================================================
- Tone: {tone_instruction}
- Target Audience: {audience_instruction}
- Do not add ```markdown at start and ``` at the end of response.
- Avoid generic phrasing like "in today’s fast-paced world" or "businesses need to adapt."
- Use **specific**, **action-oriented**, and **emotive** language.
- Maintain storytelling rhythm: short lines that sound natural as voiceover.
- Keep the flow: Hook → Insight → Value → CTA.
- visual and narration heading should be in bold.
- start visual and narration in new line
- Generate output in markdown format

=====================================================
📚 CONTEXT & REFERENCE
=====================================================
Industry: {industry or "Not specified"}
Topic: "{query}"
Reference Content:
{final_content or "[No reference material provided]"}

=====================================================
💡 OUTPUT FORMAT (Example)
=====================================================
Return only the **final timestamped script** like this:

0:00–0:{scene_duration:02d} → [Scene 1: Problem introduction narration and visuals]  
0:{scene_duration:02d}–0:{scene_duration*2:02d} → [Scene 2: Brand introduction narration and visuals]  

"""

# CTA mapping

cta_mapping = {
    "Book Assessment": "Book your free SAP Clean Core Assessment today.",
    "Request a demo": "Request a demo to see the solution in action.",
    "Talk to our experts": "Talk to our experts to discuss your requirements.",
    "Learn more about our solutions": "Learn more about our solutions tailored to your needs.",
    "Contact us today": "Contact us today to get started.",
    "Download the full guide": "Download the full guide to explore more insights.",
    "Book a free consultation": "Book your free consultation today.",
}

# Generate / Refine Handlers

def call_openai_chat(client, prompt_system, max_tokens=3200, temperature=0.7):
    """Call Azure OpenAI chat completion and return content or raise"""
    messages = [{"role": "system", "content": prompt_system}]
    response = client.chat.completions.create(messages=messages, model="Codetest", max_tokens=max_tokens, temperature=temperature)
    return response.choices[0].message.content


if generate_button and query:
    # Validate SEO for blog
    if content_type == "Blog" and not primary_keyword.strip():
        st.error("For Blog generation, Primary Keyword is required.")
    else:
        with st.spinner(f"Generating {content_type}..."):
            try:
                db, client = init_services()
                if not db or not client:
                    raise Exception("Failed to initialize one or more services.")
            except Exception as e:
                st.error(f"Service init error: {e}")
                db, client = None, None
                st.stop()  # Stop execution if services fail

            # Build full query (topic + additional)
            full_query = query.strip()

            final_content = retrieve_content(full_query, uploaded_files, url_list, db)

            # Prepare CTA text
            cta_text = cta_mapping.get(cta_choice, cta_choice)

            # Build prompt
            if content_type == "Blog":
                prompt = generate_blog_prompt(tone, target_audience, industry, full_query, word_limit, final_content,
                                              primary_keyword.strip(), lsi_keywords, cta_text)
            else:
                prompt = generate_video_prompt(tone, target_audience, industry, final_content, cta_text, query, time_limit)

            try:
                output = call_openai_chat(client, prompt)
            except Exception as e:
                st.error(f"OpenAI API error: {e}")
                output = ""

            if output:
                st.session_state.output = output
                st.session_state.last_prompt = full_query

                # Run SEO checks if blog
                if content_type == "Blog":
                    #seo_results = seo_check(output, primary_keyword.strip(), lsi_keywords)
                    st.session_state.seo_results = {}

                # Scroll or re-render to show output
                st.rerun()


# Apply refine
if apply_refine and st.session_state.output and refine_instruction and refine_instruction.strip():
    with st.spinner("Applying refinement..."):
        try:
            db, client = init_services()
            if not db or not client:
                raise Exception("Failed to initialize one or more services.")
        except Exception as e:
            st.error(f"Service init error: {e}")
            db, client = None, None
            st.stop()  # Stop execution if services fail

        refine_prompt = f"Refine the following content based on instruction: '{refine_instruction.strip()}'\n\nContent:\n{st.session_state.output}"
        try:
            new_output = call_openai_chat(client, refine_prompt, max_tokens=3000)
        except Exception as e:
            st.error(f"Refinement error: {e}")
            new_output = ""

        if new_output:
            st.session_state.output = new_output
            # Re-run SEO validation if blog
            if content_type == "Blog":
                #seo_results = seo_check(new_output, primary_keyword.strip(), lsi_keywords)
                st.session_state.seo_results = {}
            st.rerun()

# End of app

